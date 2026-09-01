#!/usr/bin/env python3
"""Assemble docs/deck-export/*.png into a 16:9 .pptx, one slide per PNG.

Each slide is the rendered PNG placed full-bleed, with the slide's visible text
copied into the speaker notes so the deck stays searchable and presentable.

The deck's look is carried by two webfonts and a lot of CSS, so a text-native
pptx would lose it. Shipping images keeps the design exactly as reviewed; the
notes are what make the file useful rather than just pretty.

Usage: python3 scripts/build-pptx.py [--out docs/Valentin-deck.pptx]
Run scripts/deck-slide-shots.ts first to produce the PNGs.
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

EXPORT = Path("docs/deck-export")
DEFAULT_OUT = Path("docs/Valentin-deck.pptx")

# 16:9 at the size PowerPoint treats as widescreen default.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def build(out: Path) -> int:
    manifest = EXPORT / "slides.json"
    if not manifest.exists():
        sys.exit(
            f"{manifest} not found — run `npx tsx scripts/deck-slide-shots.ts` first."
        )

    data = json.loads(manifest.read_text())
    slides = data["slides"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # the only layout with no placeholders to fight

    for entry in slides:
        png = EXPORT / entry["file"]
        if not png.exists():
            sys.exit(f"missing {png} — re-run the shot script.")

        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), Emu(0), Emu(0), width=SLIDE_W, height=SLIDE_H)

        notes = entry.get("notes", "").strip()
        if entry.get("clipped"):
            notes = f"[render warning: slide overflowed by {entry['clipped']}]\n\n{notes}"
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    size_mb = out.stat().st_size / 1_000_000
    print(f"{len(slides)} slides → {out} ({size_mb:.1f} MB)")

    clipped = [s for s in slides if s.get("clipped")]
    if clipped:
        print(f"\n{len(clipped)} slide(s) were clipped when rendered:")
        for s in clipped:
            print(f"  {s['file']}: {s['clipped']}")
    return 0


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT)
    sys.exit(build(ap.parse_args().out))
